import streamlit as st
import pandas as pd
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
import os
import html
import markdown

# 定义应用的标题
st.title("多文件类型阅读器")

# 创建一个文件上传器组件
uploaded_files = st.file_uploader("上传文件", type=["txt", "pdf", "docx", "pptx", "xlsx", "html", "md", "csv"], accept_multiple_files=True)

# 检查是否有文件被上传
if uploaded_files:
    if uploaded_files:
        for uploaded_file in uploaded_files:
            # 获取文件名
            file_name = uploaded_file.name
            file_type = file_name.split('.')[-1].lower()

            # 根据文件类型处理文件内容
            if file_type == 'txt':
                content = uploaded_file.getvalue().decode('utf-8')  # Use getvalue() instead of read()
                st.subheader(f"{file_name} (TXT)")
                st.text(content)

            elif file_type == 'pdf':
                pdf_reader = PdfReader(uploaded_file)  # 使用新的 PdfReader 类
                num_pages = len(pdf_reader.pages)
                st.subheader(f"{file_name} (PDF)")
                st.write(f"页数: {num_pages}")
                # 读取PDF内容
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()  # 使用新的 extract_text() 方法
                    st.write(f"第 {page_num + 1} 页内容:")
                    st.text(text)


            elif file_type == 'docx':
                doc = Document(uploaded_file)
                content = '\n'.join([para.text for para in doc.paragraphs])
                st.subheader(f"{file_name} (DOCX)")
                st.text(content)

            elif file_type == 'pptx':
                ppt = Presentation(uploaded_file)
                content = '\n'.join([slide.shapes.title.text for slide in ppt.slides if slide.shapes.title])
                st.subheader(f"{file_name} (PPTX)")
                st.text(content)

            elif file_type == 'xlsx':
                data = pd.read_excel(uploaded_file)
                st.subheader(f"{file_name} (XLSX)")
                st.dataframe(data)

            elif file_type == 'html':
                content = uploaded_file.getvalue().decode('utf-8')  # Use getvalue() instead of read()
                st.subheader(f"{file_name} (HTML)")
                st.markdown(content, unsafe_allow_html=True)

            elif file_type == 'md':
                content = uploaded_file.getvalue().decode('utf-8')  # Use getvalue() instead of read()
                st.subheader(f"{file_name} (Markdown)")
                st.markdown(markdown.markdown(content), unsafe_allow_html=True)

            elif file_type == 'csv':
                data = pd.read_csv(uploaded_file)
                st.subheader(f"{file_name} (CSV)")
                st.dataframe(data)

# 创建一个按钮，点击后会显示一条消息
if st.button("查看效果"):
    st.write("这里是上传文件读取效果！")



# --------------------

# 定义应用的标题
st.title("多文件类型上传并保存在知识库内")

# 定义保存文件的目录
script_dir = os.path.dirname(os.path.abspath(__file__))

save_directory = os.path.join(script_dir, 'docs')

# 确保保存目录存在
if not os.path.exists(save_directory):
    os.makedirs(save_directory)

# 创建一个文件上传器组件
uploaded_files_to_save = st.file_uploader("上传文件保存", type=["txt", "pdf", "docx", "pptx", "xlsx", "html", "md", "csv"], accept_multiple_files=True, key="file_uploader_to_save")

# 检查是否有文件被上传
if uploaded_files_to_save:
    for uploaded_file in uploaded_files_to_save:
        # 获取文件名
        file_name = uploaded_file.name
        file_type = file_name.split('.')[-1].lower()

        # 根据文件类型处理文件内容
        if file_type == 'txt':
            content = uploaded_file.getvalue().decode('utf-8')
            st.subheader(f"{file_name} (TXT)")
            st.text(content)
            # 保存TXT文件
            save_path = os.path.join(save_directory, file_name)
            with open(save_path, "w", encoding='utf-8') as f:
                f.write(content)

        elif file_type == 'pdf':
            pdf_reader = PdfReader(uploaded_file)
            num_pages = len(pdf_reader.pages)
            st.subheader(f"{file_name} (PDF)")
            st.write(f"页数: {num_pages}")
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                st.write(f"第 {page_num + 1} 页内容:")
                st.text(text)
                # 保存PDF文件（这里需要一个PDF写入库，如PyFPDF或reportlab）
                # 注意：PyPDF2不支持直接保存修改后的PDF内容

        elif file_type == 'docx':
            doc = Document(uploaded_file)
            content = '\n'.join([para.text for para in doc.paragraphs])
            st.subheader(f"{file_name} (DOCX)")
            st.text(content)
            # 保存DOCX文件
            doc.save(os.path.join(save_directory, file_name))

        elif file_type == 'pptx':
            ppt = Presentation(uploaded_file)
            content = '\n'.join([slide.shapes.title.text for slide in ppt.slides if slide.shapes.title])
            st.subheader(f"{file_name} (PPTX)")
            st.text(content)
            # 保存PPTX文件
            ppt.save(os.path.join(save_directory, file_name))

        elif file_type == 'xlsx':
            data = pd.read_excel(uploaded_file)
            st.subheader(f"{file_name} (XLSX)")
            st.dataframe(data)
            # 保存XLSX文件
            data.to_excel(os.path.join(save_directory, file_name), index=False)

        elif file_type in ['html', 'md']:
            content = uploaded_file.getvalue().decode('utf-8')
            st.subheader(f"{file_name} ({file_type})")
            if file_type == 'html':

                st.write(content,unsafe_allow_html=True)
            else:
                st.markdown(content, unsafe_allow_html=True)
            # 保存HTML和MD文件
            save_path = os.path.join(save_directory, file_name)
            with open(save_path, "w", encoding='utf-8') as f:
                f.write(content)

        elif file_type == 'csv':
            data = pd.read_csv(uploaded_file)
            st.subheader(f"{file_name} (CSV)")
            st.dataframe(data)
            # 保存CSV文件
            data.to_csv(os.path.join(save_directory, file_name), index=False)

# 创建一个按钮，点击后会显示一条消息
if st.button("显示消息"):
    st.write("感谢上传文件！")

